from __future__ import annotations

import re
from pathlib import Path, PurePath, PureWindowsPath

from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document
from pptx import Presentation

from exam_helper.config import SUPPORTED_EXTENSIONS


def _safe_upload_name(raw_name: str, fallback_index: int, used_names: set[str]) -> str:
    base_name = PurePath(PureWindowsPath(str(raw_name)).name).name
    sanitized = re.sub(r"[^A-Za-z0-9._ -]+", "_", base_name).strip(" .")
    if not sanitized:
        sanitized = f"upload-{fallback_index}"

    candidate_path = Path(sanitized)
    suffix = candidate_path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"Unsupported file type for {raw_name!r}. Supported types: {supported}")

    stem = candidate_path.stem.strip(" .") or f"upload-{fallback_index}"
    candidate = f"{stem}{suffix}"
    counter = 2
    while candidate.lower() in used_names:
        candidate = f"{stem}-{counter}{suffix}"
        counter += 1

    used_names.add(candidate.lower())
    return candidate


def save_uploaded_files(uploaded_files: list, module_dir: Path) -> list[Path]:
    uploads_dir = module_dir / "uploads"
    uploads_dir.mkdir(parents=True, exist_ok=True)

    file_paths: list[Path] = []
    used_names: set[str] = set()
    for idx, uploaded in enumerate(uploaded_files, start=1):
        file_name = _safe_upload_name(uploaded.name, idx, used_names)
        path = uploads_dir / file_name
        path.write_bytes(uploaded.getvalue())
        file_paths.append(path)
    return file_paths


def _normalize_document(text: str, source_file: str, source_type: str, page_or_slide_number: int) -> Document | None:
    stripped = text.strip()
    if not stripped:
        return None
    return Document(
        page_content=stripped,
        metadata={
            "source_file": source_file,
            "source_type": source_type,
            "page_or_slide_number": int(page_or_slide_number),
        },
    )


def load_pdf(path: Path) -> list[Document]:
    loader = PyPDFLoader(str(path))
    raw_docs = loader.load()
    docs: list[Document] = []

    for raw in raw_docs:
        page_idx = raw.metadata.get("page", raw.metadata.get("page_number", 0))
        normalized = _normalize_document(
            text=raw.page_content,
            source_file=path.name,
            source_type="pdf",
            page_or_slide_number=int(page_idx) + 1,
        )
        if normalized:
            docs.append(normalized)
    return docs


def load_pptx(path: Path) -> list[Document]:
    presentation = Presentation(str(path))
    docs: list[Document] = []
    for i, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        normalized = _normalize_document(
            text="\n".join(texts),
            source_file=path.name,
            source_type="pptx",
            page_or_slide_number=i,
        )
        if normalized:
            docs.append(normalized)
    return docs


def load_documents_from_paths(paths: list[Path]) -> list[Document]:
    docs: list[Document] = []
    for path in paths:
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            docs.extend(load_pdf(path))
        elif suffix == ".pptx":
            docs.extend(load_pptx(path))
        else:
            raise ValueError(f"Unsupported file type: {path.name}")
    return docs
