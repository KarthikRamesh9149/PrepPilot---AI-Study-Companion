from pptx import Presentation

from exam_helper.ingestion.loaders import load_pptx


def test_loaders_pptx_metadata(tmp_path):
    pptx_path = tmp_path / "sample.pptx"
    presentation = Presentation()

    slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide1.shapes.title.text = "Slide One"
    slide1.placeholders[1].text = "Intro content"

    slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide2.shapes.title.text = "Slide Two"
    slide2.placeholders[1].text = "More content"

    presentation.save(str(pptx_path))

    docs = load_pptx(pptx_path)
    assert len(docs) == 2
    assert docs[0].metadata["source_file"] == "sample.pptx"
    assert docs[0].metadata["source_type"] == "pptx"
    assert docs[0].metadata["page_or_slide_number"] == 1
    assert docs[1].metadata["page_or_slide_number"] == 2
